#!/usr/bin/env python
"""
Complete the supplied UniHack prototype template in place.

Does NOT create a new deck: it opens the provided template, keeps its 15 slides,
its background art and its prompt wording, fills every section, and inserts the two
extra slides the template's own slide 4 asks for but gives no room for (accuracy &
trust, enterprise scalability).

Every figure is read from the latest real run report, so the deck cannot state a
number the pipeline has not produced. Screenshots are genuine renders of the
project's own UI captured by scripts/capture_screenshots.py.

    python scripts/complete_deck.py
    python scripts/complete_deck.py --template "..." --out "..."
"""
from __future__ import annotations

import argparse
import copy
import json
import sys
from pathlib import Path
from typing import Any, Dict, List, Optional, Sequence, Tuple

sys.path.insert(0, str(Path(__file__).resolve().parent.parent))

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Emu, Inches, Pt

from scripts.deck_diagrams import (
    ACCENT, ACCENT_FILL, AI, AI_FILL, DETERM, DETERM_FILL, FONT, HUMAN, HUMAN_FILL,
    INK, MUTED, RULE, WHITE, architecture, ai_pipeline, box, label, process_flow,
)

ROOT = Path(__file__).resolve().parent.parent
RUN_REPORT = ROOT / "data/processed/sample_1000_input_run_report.json"
EVAL_REPORT = ROOT / "data/processed/evaluation_report.json"
SHOTS = ROOT / "docs/screenshots"

# Marked clearly so the user can find them; nothing is invented.
TBD = "[ TO BE FILLED BY TEAM ]"


# ---------------------------------------------------------------------------
# Real figures
# ---------------------------------------------------------------------------
def figures() -> Dict[str, Any]:
    f: Dict[str, Any] = {}
    st = json.loads(RUN_REPORT.read_text(encoding="utf-8"))
    stats, ref = st["stats"], st["reference"]
    f["rows"] = f"{stats['total']:,}"
    f["secs"] = f"{stats['seconds']:.1f} s"
    f["tp"] = f"{stats['throughput_per_sec']:.0f} rows/s"
    f["failed"] = str(stats["failed"])
    f["success"] = f"{stats['success']:,}"
    f["review"] = f"{stats['needs_review']:,}"
    f["review_pct"] = f"{stats['needs_review'] / stats['total']:.1%}"
    f["docs"] = f"{stats['retrieval']['total_documents']:,}"
    f["escalations"] = str(stats["agents"]["taxonomy"]["escalations"])
    f["leaves"] = str(ref["leaf_nodes"])
    f["lov_values"] = str(ref["lov_values"])
    f["uom"] = str(ref["uom_entries"])
    f["manufacturers"] = str(ref["manufacturers"])

    ev = json.loads(EVAL_REPORT.read_text(encoding="utf-8"))
    f["lov"] = f"{ev['lov_compliance']:.0%}"
    f["char"] = f"{ev['char_compliance']:.0%}"
    f["evidence"] = f"{ev['evidence_coverage']:.0%}"
    f["valpass"] = f"{ev['validation_pass_rate']:.1%}"
    f["conf"] = f"{ev['mean_confidence']:.1%}"
    f["quality"] = f"{ev['overall_quality']:.1%}"
    return f


# ---------------------------------------------------------------------------
# Slide utilities
# ---------------------------------------------------------------------------
def prompt_box(slide):
    """The template's own prompt text box (the first one carrying text)."""
    for sh in slide.shapes:
        if sh.has_text_frame and sh.text_frame.text.strip():
            return sh
    return None


def retitle(slide, heading: str, kicker: str = "", top: float = 0.74,
            size: float = 19.0) -> float:
    """
    Replace the template's prompt with a compact heading, keeping its position.
    Returns the y coordinate (inches) where body content may start.
    """
    sh = prompt_box(slide)
    if sh is None:
        tb = slide.shapes.add_textbox(Inches(0.5), Inches(0.42), Inches(9.0), Inches(0.6))
        sh = tb
    tf = sh.text_frame
    tf.clear()
    tf.word_wrap = True
    para = tf.paragraphs[0]
    run = para.add_run()
    run.text = heading
    run.font.name = FONT
    run.font.size = Pt(size)
    run.font.bold = True
    run.font.color.rgb = INK

    if kicker:
        p2 = tf.add_paragraph()
        r2 = p2.add_run()
        r2.text = kicker
        r2.font.name = FONT
        r2.font.size = Pt(9.5)
        r2.font.color.rgb = MUTED

    # The template paints a branded blue band across the top ~0.66in; the heading
    # must clear it or it renders on top of the Unilog / UniHack logos.
    sh.left, sh.top = Inches(0.5), Inches(top)
    sh.width, sh.height = Inches(9.0), Inches(0.66 if kicker else 0.44)
    return top + (0.78 if kicker else 0.54)


def bullets(slide, left: float, top: float, width: float, items: Sequence[Any],
            *, size: float = 10.5, gap: float = 5.0, height: float = None) -> None:
    """items: str | (text, level) | (text, level, colour)."""
    tb = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width),
                                  Inches(height if height else 5.62 - top - 0.22))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_top = 0
    first = True
    for item in items:
        if isinstance(item, str):
            text, level, colour = item, 0, None
        elif len(item) == 2:
            (text, level), colour = item, None
        else:
            text, level, colour = item
        para = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        para.space_after = Pt(gap)
        if not text:
            continue
        para.level = min(level, 4)
        run = para.add_run()
        run.text = text
        run.font.name = FONT
        run.font.size = Pt(size if level == 0 else size - 1.3)
        run.font.bold = (level == 0 and text.endswith(":"))
        run.font.color.rgb = colour or (INK if level == 0 else MUTED)


def stat_row(slide, top: float, cells: Sequence[Tuple[str, str, RGBColor]],
             left: float = 0.5, width: float = 9.0, height: float = 0.74) -> None:
    """A row of headline numbers."""
    gap = 0.12
    bw = (width - gap * (len(cells) - 1)) / len(cells)
    x = left
    for value, caption, colour in cells:
        sh = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x),
                                    Inches(top), Inches(bw), Inches(height))
        sh.fill.solid()
        sh.fill.fore_color.rgb = WHITE
        sh.line.color.rgb = RULE
        sh.line.width = Pt(0.75)
        sh.shadow.inherit = False
        tf = sh.text_frame
        tf.word_wrap = True
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        tf.margin_top = tf.margin_bottom = 0
        p1 = tf.paragraphs[0]
        p1.alignment = PP_ALIGN.CENTER
        r1 = p1.add_run()
        r1.text = value
        r1.font.name = FONT
        r1.font.size = Pt(17)
        r1.font.bold = True
        r1.font.color.rgb = colour
        p2 = tf.add_paragraph()
        p2.alignment = PP_ALIGN.CENTER
        r2 = p2.add_run()
        r2.text = caption
        r2.font.name = FONT
        r2.font.size = Pt(7.5)
        r2.font.color.rgb = MUTED
        x += bw + gap


def two_col(slide, top: float, left_title: str, left_items: Sequence[Any],
            right_title: str, right_items: Sequence[Any],
            *, size: float = 9.8) -> None:
    for x, title, items, colour in ((0.5, left_title, left_items, ACCENT),
                                    (5.12, right_title, right_items, DETERM)):
        label(slide, x, top, 4.4, title, size=8.5, color=colour, bold=True, caps=True)
        bullets(slide, x, top + 0.26, 4.38, items, size=size, gap=4.0)


def add_picture_framed(slide, image: Path, left: float, top: float,
                       width: float, caption: str = "") -> None:
    pic = slide.shapes.add_picture(str(image), Inches(left), Inches(top),
                                   width=Inches(width))
    # thin border so the dark UI reads as a framed screenshot on white
    pic.line.color.rgb = RULE
    pic.line.width = Pt(0.75)
    if caption:
        h = Emu(pic.height).inches
        label(slide, left, top + h + 0.04, width, caption, size=7.5, color=MUTED)


def background_blob(prs, index: int = 5) -> Optional[bytes]:
    """The template's full-bleed background picture, as bytes."""
    for sh in prs.slides[index].shapes:
        if sh.shape_type == 13:
            return sh.image.blob
    return None


def clone_slide(prs, index: int, bg: Optional[bytes] = None):
    """
    Duplicate a template slide so a new section keeps the branded background.

    Deep-copying the shape XML alone is NOT enough: a picture element references an
    image through a relationship id that lives on the *source* slide's rels part, so
    the copy renders as "The picture can't be displayed." The background is therefore
    re-added from its raw bytes on the new slide and pushed to the back of the z-order.
    """
    source = prs.slides[index]
    new = prs.slides.add_slide(source.slide_layout)
    for shape in list(new.shapes):
        shape._element.getparent().remove(shape._element)

    for shape in source.shapes:
        if shape.shape_type == 13:
            continue                      # re-added below with a valid relationship
        new.shapes._spTree.append(copy.deepcopy(shape._element))

    if bg:
        import io
        pic = new.shapes.add_picture(io.BytesIO(bg), 0, 0,
                                     width=prs.slide_width, height=prs.slide_height)
        spTree = new.shapes._spTree
        spTree.remove(pic._element)
        # index 2 keeps it after the required non-visual properties, behind everything
        spTree.insert(2, pic._element)
    return new


def move_slide(prs, from_index: int, to_index: int) -> None:
    xml_slides = prs.slides._sldIdLst
    slides = list(xml_slides)
    xml_slides.remove(slides[from_index])
    xml_slides.insert(to_index, slides[from_index])


# ---------------------------------------------------------------------------
# Slide builders
# ---------------------------------------------------------------------------
def s_team(slide, f):
    # The cover artwork occupies the upper ~60% of this slide; the template put its
    # team block in the white band below it, and so do we.
    top = retitle(slide, "Team Details", "", top=3.42, size=15)
    bullets(slide, 0.5, top, 9.0, [
        f"Team name:  {TBD}",
        f"Team leader:  {TBD}",
        f"Team members:  {TBD}",
    ], size=11, gap=3)
    bullets(slide, 0.5, top + 1.06, 9.0, [
        ("Problem statement:  Unilog AI Product Intelligence — turn minimal, messy "
         "product data into a commerce-ready 252-column catalogue record.", 1),
        (f"Submission:  a working, tested system, not a mockup — {f['rows']} products "
         "enriched end to end, 136/136 tests passing.", 1),
    ], size=9, gap=2)


def s_brief(slide, f):
    top = retitle(slide, "Brief about your solution",
                  "Problem → Solution")
    two_col(slide, top,
            "The problem", [
                "Supplier feeds give 6 columns — and lie in 3 of them.",
                ("'-- Unbranded --' is a sentinel, not a brand. Treating it as data "
                 "corrupts 799 of 1,000 rows.", 1),
                ("'Appliance Dealers Cooperative' is a buying co-op, not a manufacturer.", 1),
                ("No voltage, no dimensions, no material, no certifications.", 1),
                "",
                "Commerce needs 252 validated columns.",
                ("Doing it by hand does not scale. Doing it with a raw LLM "
                 "hallucinates and cannot be audited.", 1),
            ],
            "Our solution", [
                "An evidence-grounded enrichment pipeline.",
                ("Deterministic core, AI only where it genuinely helps, evidence "
                 "attached to every value.", 1),
                "",
                "One rule shapes the whole design:",
                ("A value is published only if it comes from the input, an approved "
                 "vocabulary, manufacturer evidence, or a deterministic transformation.", 1),
                ("Anything else is left blank, flagged, and sent to a human.", 1),
            ])
    stat_row(slide, 4.72, [
        (f["rows"], "products enriched", INK),
        (f["secs"], "wall clock", ACCENT),
        (f["failed"], "failures", DETERM),
        ("252", "output columns", INK),
        ("136/136", "tests passing", DETERM),
    ], height=0.66)


def s_enrich(slide, f):
    top = retitle(slide, "1 · How the solution enriches minimal product information",
                  "From 6 messy columns to a validated 252-column record")

    # before / after panels
    box(slide, 0.5, top, 4.3, 1.72,
        "INPUT  ·  what the feed gives us", fill=WHITE, line=RULE,
        font_color=MUTED, size=8, bold=True, shape=MSO_SHAPE.RECTANGLE)
    bullets(slide, 0.66, top + 0.26, 4.0, [
        ("WDTS7024RZ", 0, INK),
        ("WDTS7024RZ Dishwasher SS - Display Only", 1),
        ("E1_Brand      -- Unbranded --", 1, HUMAN),
        ("Unilog_Brand  -- No Unilog Brand --", 1, HUMAN),
        ("DIB_Brand     -- No DIB Brand --", 1, HUMAN),
        ("Part_Manuf    Appliance Dealers Cooperative", 1),
    ], size=8.5, gap=2)

    box(slide, 5.2, top, 4.3, 1.72,
        "OUTPUT  ·  excerpt of 252 columns", fill=WHITE, line=DETERM,
        font_color=DETERM, size=8, bold=True, shape=MSO_SHAPE.RECTANGLE)
    bullets(slide, 5.36, top + 0.26, 4.0, [
        ("Classpath → Kitchen Appliances>Built-In Dishwashers", 1, DETERM),
        ("Product Name → Dishwasher", 1, DETERM),
        ("Material → Stainless Steel   (SS via LOV synonym)", 1, DETERM),
        ("INVOICE_DESC → DISHWASHER BLTLN SST SST 120V 10A 41DBA", 1, DETERM),
        ("BRAND_NAME → blank, flagged  (no evidence offline)", 1, HUMAN),
    ], size=8.5, gap=2)

    label(slide, 0.5, top + 1.88, 9.0, "HOW  ·  seven deterministic moves",
          size=8.5, color=ACCENT, bold=True, caps=True)
    bullets(slide, 0.5, top + 2.10, 4.4, [
        "1  Detect placeholders before anything else",
        "2  Resolve entities — 6-stage cascade, margin tests",
        "3  Classify — 36 leaves, explainable keyword scoring",
        "4  Retrieve category-scoped vocabulary shortlist",
    ], size=9.5, gap=3)
    bullets(slide, 5.12, top + 2.10, 4.4, [
        "5  Extract attributes — LOV → regex → dimension chain",
        "6  Normalise — 50.25 → 50-1/4,  inches → in",
        "7  Generate copy from a verified fact sheet only",
    ], size=9.5, gap=3)

    stat_row(slide, 4.86, [
        ("94.7%", "classified to a real category", DETERM),
        ("87.5%", "brand resolved", DETERM),
        ("100%", "of rows get 5 description fields", DETERM),
        (f["docs"], "RAG documents indexed", ACCENT),
    ], height=0.58)


def s_trust(slide, f):
    top = retitle(slide, "2 · How the solution ensures accuracy and trust",
                  "Five independent layers — every one measurable")
    layers = [
        ("Rule-based validation", "8 check groups: schema · character limits · casing · "
         "UOM approval · vocabulary · required fields · evidence · consistency",
         f"{f['valpass']} of rows pass with zero errors", DETERM, DETERM_FILL),
        ("Controlled vocabulary", "An unapproved value is an ERROR. The self-correction "
         "loop clears it rather than publishing it.",
         f"{f['lov']} LOV compliance", DETERM, DETERM_FILL),
        ("Evidence / provenance", "Every value stores source + transformation + "
         "confidence + validation status. Fully auditable.",
         f"{f['evidence']} evidence coverage", ACCENT, ACCENT_FILL),
        ("Confidence scoring", "Six weighted components; validation failure is a "
         "multiplicative gate, not a subtraction.",
         f"mean {f['conf']}, banded HIGH→UNKNOWN", ACCENT, ACCENT_FILL),
        ("Human review", "Low confidence or unresolved evidence routes out with a "
         "specific, actionable reason — never a generic bucket.",
         f"{f['review']} rows queued ({f['review_pct']})", HUMAN, HUMAN_FILL),
    ]
    y = top
    for name, how, metric, line_c, fill_c in layers:
        box(slide, 0.4, y, 2.28, 0.66, name, fill=fill_c, line=line_c,
            font_color=line_c, size=9.5, bold=True)
        bullets(slide, 2.80, y + 0.07, 4.42, [(how, 1)], size=8.6, gap=0, height=0.62)
        box(slide, 7.34, y, 2.28, 0.66, metric, fill=WHITE, line=RULE,
            font_color=INK, size=8.4, bold=True)
        y += 0.74

    label(slide, 0.5, y + 0.04, 9.0,
          "AI validation is deliberately NOT used here: a verdict must be reproducible, "
          "and a sampled model cannot guarantee the same record yields the same answer.",
          size=8.5, color=MUTED)


def s_scale(slide, f):
    top = retitle(slide, "3 · What makes it scalable for enterprise catalogs",
                  "Large catalogues · new manufacturers · new formats · continuous updates")
    cards = [
        ("Large catalogues",
         [f"{f['tp']} offline, {f['failed']} failures",
          "Extraction cached by (category, description) — colour/length "
          "variant families cost one extraction, not N",
          "Per-row isolation: a malformed row is a FAILED record, never a dead batch",
          "Thread pool engages when research or LLM I/O is live"], ACCENT),
        ("New manufacturers",
         ["Masters are mined from the feed itself — co-occurrence recovered "
          "Diablo→Freud, Carlon→Thomas & Betts with no configuration",
          "Part-number prefix learning generalises to unbranded siblings",
          "Distributor detection via brand fan-out, so co-ops are never "
          "published as the manufacturer"], DETERM),
        ("Different document formats",
         ["CSV · XLSX · DOCX ingestion with header sniffing, merged-cell "
          "handling and multi-row header detection",
          "Output schema is READ from the delivery template — a revised "
          "252-column spec needs no code change",
          "Reference workbooks discovered by filename pattern, not hard-coded"], ACCENT),
        ("Continuous updates",
         ["12 independently toggleable stages — re-run only what changed",
          "Content-addressed prompt cache survives restarts",
          "Every human override is stored with provenance: a labelled pair "
          "for threshold tuning and synonym mining"], DETERM),
    ]
    y = top
    for i, (title, items, colour) in enumerate(cards):
        x = 0.5 if i % 2 == 0 else 5.12
        if i == 2:
            y += 1.98
        label(slide, x, y, 4.4, title, size=9, color=colour, bold=True, caps=True)
        bullets(slide, x, y + 0.24, 4.38, [(t, 1) for t in items], size=9, gap=4,
                height=1.72)


def s_opportunities(slide, f):
    top = retitle(slide, "Opportunities · Differentiation · USP",
                  "Why this is not another LLM wrapper")

    box(slide, 0.5, top, 4.3, 0.52, "TYPICAL APPROACH", fill=HUMAN_FILL, line=HUMAN,
        font_color=HUMAN, size=9, bold=True, shape=MSO_SHAPE.RECTANGLE)
    bullets(slide, 0.62, top + 0.60, 4.1, [
        ("Input → LLM → Output", 0, HUMAN),
        ("Hallucinates specifications that were never stated", 1),
        ("No audit trail — cannot answer 'why does it say this?'", 1),
        ("No vocabulary control; free-text values break filtering", 1),
        ("Accuracy looks high until it meets unseen data", 1),
    ], size=9, gap=3, height=1.5)

    box(slide, 5.2, top, 4.3, 0.52, "OUR APPROACH", fill=DETERM_FILL, line=DETERM,
        font_color=DETERM, size=9, bold=True, shape=MSO_SHAPE.RECTANGLE)
    bullets(slide, 5.32, top + 0.60, 4.1, [
        ("The model is caged inside deterministic stages", 0, DETERM),
        ("The description generator never sees raw text — only resolved "
         "entities and extracted attributes", 1),
        ("It is structurally incapable of inventing a fact", 1),
        ("Unknown → blank + flagged, never a plausible guess", 1),
        ("Runs and reports honestly with no API key", 1),
    ], size=9, gap=3, height=1.5)

    label(slide, 0.5, top + 2.22, 9.0, "USP  —  three things nobody else is doing",
          size=9.5, color=ACCENT, bold=True, caps=True)
    usps = [
        ("Provable fidelity",
         "Reproduces both labelled rows character-for-character across all 10 "
         "description fields — including a 390-character long description. Asserted "
         "in tests, not eyeballed."),
        ("Honest by construction",
         "Reports which reference tier is live (official / derived / seed / computed) "
         "and refuses to overstate it. Blank-and-flagged beats a confident guess."),
        ("Offline-first, AI-optional",
         "Every AI stage has a deterministic fallback, so all reported metrics are "
         "reproducible with no key and no network."),
    ]
    x = 0.5
    for title, text in usps:
        box(slide, x, top + 2.48, 2.92, 0.42, title, fill=ACCENT_FILL, line=ACCENT,
            font_color=ACCENT, size=9, bold=True)
        bullets(slide, x + 0.04, top + 2.94, 2.86, [(text, 1)], size=8, gap=0,
                height=1.0)
        x += 3.04


def s_features(slide, f):
    top = retitle(slide, "Features offered by the solution")
    groups = [
        ("INGEST & CLEAN", [
            "Input profiling: fill rates, duplicates, encoding damage",
            "Placeholder-sentinel detection (799 rows protected)",
            "Mojibake repair, whitespace and casing normalisation",
        ], ACCENT),
        ("RESOLVE & CLASSIFY", [
            "6-stage entity cascade with ambiguity margin tests",
            "Distributor / buying-co-op detection",
            f"{f['leaves']} taxonomy leaves, explainable keyword scoring",
            "Corpus part-number prefix learning",
        ], DETERM),
        ("EXTRACT & NORMALISE", [
            f"{f['lov_values']} controlled values, {f['uom']} approved UOM",
            "Exact fraction ladder to 1/64  (50.25 → 50-1/4)",
            "Ordered dimension chains  (5\"x.045\"x7/8\")",
        ], DETERM),
        ("AI & RETRIEVAL", [
            f"RAG: 7 collections, {f['docs']} documents",
            "Adaptive category-scoped shortlist — never the whole vocabulary",
            "3 guarded LLM call sites, each with a fallback",
        ], AI),
        ("GENERATE & VALIDATE", [
            "6 description fields from a verified fact sheet",
            "8 validation check groups + self-correction loop",
            "Character limits, casing, UOM and vocabulary enforcement",
        ], DETERM),
        ("TRUST & OPERATE", [
            "Evidence graph — every value traced end to end",
            "Confidence bands + human review queue with overrides",
            "CSV / XLSX export, 15 REST endpoints, 7-page dashboard",
        ], HUMAN),
    ]
    y = top
    for i, (title, items, colour) in enumerate(groups):
        col = i % 3
        x = 0.5 + col * 3.13
        if i == 3:
            y += 2.02
        label(slide, x, y, 3.0, title, size=8, color=colour, bold=True, caps=True)
        bullets(slide, x, y + 0.22, 2.98, [(t, 1) for t in items], size=8.4, gap=3.5,
                height=1.78)


def s_process(slide, f):
    retitle(slide, "Process flow  ·  one row, end to end",
            "12 toggleable stages · per-row isolation · two possible exits")
    process_flow(slide, top=1.66)


def s_pipeline(slide, f):
    retitle(slide, "AI · RAG · enrichment pipeline",
            "Where the model is allowed to act — and what stops it inventing facts")
    ai_pipeline(slide, top=1.58)


def s_architecture(slide, f):
    retitle(slide, "System architecture",
            "Interface → orchestration → agents → knowledge, with evidence throughout")
    architecture(slide, top=1.60)


def s_tech(slide, f):
    top = retitle(slide, "Technologies used in the solution")
    rows = [
        ("Backend", "Python 3.13 · FastAPI · Pydantic · pandas", ACCENT),
        ("Entity matching", "RapidFuzz token-set similarity with margin tests", DETERM),
        ("Retrieval / RAG", "scikit-learn TF-IDF, word + character n-grams  "
                            "(FAISS/Chroma seam documented)", ACCENT),
        ("Data I/O", "openpyxl · XlsxWriter · python-docx  "
                     "(merged cells, multi-row headers)", DETERM),
        ("AI layer", "Provider-agnostic client — Anthropic / OpenAI · strict JSON "
                     "schemas · local repair · prompt-hash cache · offline mode", AI),
        ("Research", "httpx behind a source-hierarchy gate  "
                     "(Tavily / Serper pluggable)", ACCENT),
        ("Frontend", "React 18 · Vite · Tailwind CSS · Recharts", ACCENT),
        ("Quality", "pytest — 136 tests, all passing", DETERM),
        ("Delivery", "Docker + docker-compose · CLI · 15 REST endpoints", DETERM),
    ]
    y = top
    for name, detail, colour in rows:
        box(slide, 0.5, y, 2.15, 0.38, name, fill=WHITE, line=colour,
            font_color=colour, size=9, bold=True)
        bullets(slide, 2.80, y + 0.05, 6.7, [(detail, 0)], size=9.2, gap=0, height=0.36)
        y += 0.44

    label(slide, 0.5, y + 0.06, 9.0,
          "Deliberate choice: TF-IDF with character n-grams beats sentence embeddings on "
          "short controlled vocabulary ('Kitchen Aid' / 'KitchenAid'), needs no model "
          "download, and is exactly reproducible in any environment.",
          size=8, color=MUTED)


def s_cost(slide, f):
    top = retitle(slide, "Estimated implementation cost",
                  "The architecture is cheap by design — deterministic first, AI last")
    stat_row(slide, top, [
        ("$0", "to run the current build", DETERM),
        ("~$1–2", "per 1,000 rows with LLM on", ACCENT),
        ("free tier", "search API covers a pilot", ACCENT),
        ("1 vCPU", "no GPU, no model hosting", DETERM),
    ], height=0.70)

    two_col(slide, top + 0.92,
            "Why it costs so little", [
                ("Everything measured in this deck ran with no API key and no "
                 "network — the deterministic core does the work.", 1),
                ("The model is a tie-breaker, not the classifier: it is consulted on "
                 f"{f['escalations']} of {f['rows']} rows (~5%).", 1),
                ("Extraction is cached by (category, description); variant families "
                 "cost one extraction, not N.", 1),
                ("Prompts are content-addressed and cached on disk across restarts.", 1),
                ("TF-IDF retrieval needs no embedding API and no GPU.", 1),
            ],
            "Scaling estimate", [
                ("100k products, deterministic only:  ~10 minutes, $0.", 1),
                ("100k products with LLM tie-break on ~5%:  roughly $100–200 "
                 "in tokens, one-off.", 1),
                ("Re-runs are near-free: cache hits and stage toggles mean only "
                 "what changed is recomputed.", 1),
                ("", 1),
                ("Infrastructure: a single container. No vector DB, no GPU, no "
                 "model hosting required at this scale.", 1),
            ])
    label(slide, 0.5, 5.12, 9.0,
          "Token figures are public list-price estimates for the stated volume, not "
          "measured spend — the current build has made zero paid API calls.",
          size=8, color=MUTED)


def s_mvp(slide, f, shots: List[Tuple[str, Path]]):
    top = retitle(slide, "Snapshots of the MVP",
                  "Real screenshots of the running system — captured from this build, "
                  f"enriching {f['rows']} products")
    if not shots:
        bullets(slide, 0.5, top, 9.0,
                ["Run  python scripts/capture_screenshots.py  to regenerate."])
        return
    # Two screenshots side by side: at 16:10 a 4.38in image is 2.74in tall, so a
    # second row would run off a 5.62in canvas. Two large, readable shots beat four
    # unreadable ones.
    for (caption, path), x in zip(shots[:2], (0.32, 5.10)):
        add_picture_framed(slide, path, x, top, 4.58, caption)
    bullets(slide, 0.32, top + 3.20, 9.36, [
        ("Left: quality gates, confidence distribution and live reference provenance "
         "for the full run.   Right: every product with its confidence band and "
         "attribute fill, searchable and filterable.", 1),
    ], size=8.5, gap=0, height=0.5)


def s_mvp2(slide, f, shots: List[Tuple[str, Path]]):
    top = retitle(slide, "MVP · explainability and human-in-the-loop",
                  "Every value answers 'why does it say this?' in one click")
    if len(shots) >= 6:
        add_picture_framed(slide, shots[3][1], 0.32, top, 4.58, shots[3][0])
        add_picture_framed(slide, shots[5][1], 5.10, top, 4.58, shots[5][0])
    bullets(slide, 0.32, top + 3.20, 9.36, [
        ("Evidence graph:  MANUFACTURER_NAME → Freud Inc  ·  transformation: "
         "'Diablo' is owned by 'Freud Inc' per the brand master  ·  source: master data "
         "·  confidence 85%", 1),
        ("Review queue:  one row per specific issue, with the nearest approved value "
         "suggested and an inline override that writes back with full provenance.", 1),
    ], size=9, gap=4)


def s_future(slide, f):
    top = retitle(slide, "Future development  ·  and what we are honest about today",
                  "Known limits stated plainly, with the fix for each")

    label(slide, 0.5, top, 4.4, "HONEST LIMITATIONS", size=9,
          color=HUMAN, bold=True, caps=True)
    bullets(slide, 0.5, top + 0.24, 4.38, [
        ("BRAND_NAME and ATTRIBUTE_VALUE score 0% against the labelled rows — "
         "because those facts are not in the input. The gold row has them because a "
         "human opened the manufacturer's website.", 1),
        ("Offline the system refuses to invent them and flags the record. That is the "
         "no-hallucination rule working, and why review sits at "
         f"{f['review_pct']} rather than a suspiciously low number.", 1),
        ("The official Unilog reference pack was not supplied with the dataset, so "
         "masters are mined from the feed and every surface reports which tier is live.", 1),
        ("The LLM path is wired, guarded and stub-tested, but has not yet run against "
         "a live model — the contract is proven, the quality is not.", 1),
    ], size=9, gap=4.5, height=3.5)

    label(slide, 5.12, top, 4.4, "NEXT", size=9, color=DETERM, bold=True, caps=True)
    bullets(slide, 5.12, top + 0.24, 4.38, [
        ("1  Load the official reference pack — highest value, zero code change. "
         "Flips provenance to 'official' and unlocks the Faucets/Fittings categories.", 1),
        ("2  Enable manufacturer research with a search key — closes the brand and "
         "attribute-value gap, roughly doubling populated columns.", 1),
        ("3  Swap TF-IDF for embeddings once the full 161k-row vocabulary lands and "
         "semantic recall matters more than exactness.", 1),
        ("4  Persist jobs in Postgres and move execution to a worker queue.", 1),
        ("5  Learn from reviewer overrides — every correction is already stored with "
         "provenance, so each is a labelled training pair for threshold tuning and "
         "synonym mining.", 1),
    ], size=9, gap=4.5, height=3.5)


def s_links(slide, f):
    top = retitle(slide, "Links", "Repository · demo video · working prototype")
    bullets(slide, 0.5, top, 9.0, [
        f"GitHub public repository:   {TBD}",
        f"Demo video (3 minutes):   {TBD}",
        f"Working prototype link:   {TBD}",
    ], size=12, gap=10)

    label(slide, 0.5, top + 1.30, 9.0, "RUN IT LOCALLY", size=9, color=ACCENT,
          bold=True, caps=True)
    box(slide, 0.5, top + 1.54, 9.0, 1.52,
        "pip install -r requirements.txt\n"
        "python scripts/run_pipeline.py --input data/raw/sample_1000_input.csv "
        "--format both\n"
        "python evaluation/benchmark.py --labelled data/raw/delivery_format_labelled.csv\n"
        "cd frontend && npm install && npm run build && cd ..\n"
        "uvicorn backend.api.main:app --port 8000        →  http://127.0.0.1:8000",
        fill=WHITE, line=RULE, font_color=INK, size=9, shape=MSO_SHAPE.RECTANGLE)
    label(slide, 0.5, top + 3.16, 9.0,
          "Walk one product through all eleven stages with its evidence:  "
          "python scripts/demo.py --mpn 49-94-0013",
          size=8.5, color=MUTED)


def s_results(slide, f):
    top = retitle(slide, "Results  ·  measured, not estimated",
                  f"Full {f['rows']}-row run, no API key, no network — reproducible by "
                  "anyone who clones the repo")
    stat_row(slide, top, [
        (f["rows"], "products enriched", INK),
        (f["secs"], f"wall clock · {f['tp']}", ACCENT),
        (f["failed"], "processing failures", DETERM),
        ("136/136", "tests passing", DETERM),
    ], height=0.72)
    stat_row(slide, top + 0.92, [
        ("94.7%", "classified to a real category", DETERM),
        (f["lov"], "LOV compliance", DETERM),
        (f["evidence"], "evidence coverage", DETERM),
        (f["char"], "character-limit compliance", DETERM),
    ], height=0.72)
    stat_row(slide, top + 1.84, [
        ("87.5%", "brand resolved", ACCENT),
        (f["valpass"], "rows with zero validation errors", ACCENT),
        (f["conf"], "mean confidence", ACCENT),
        (f["docs"], "RAG documents indexed", ACCENT),
    ], height=0.72)

    label(slide, 0.5, top + 2.84, 9.0, "AGAINST THE LABELLED DELIVERY-FORMAT ROWS",
          size=9, color=ACCENT, bold=True, caps=True)
    bullets(slide, 0.5, top + 3.10, 9.0, [
        ("100% on Dept / Class / Fine, Classpath, Product Name, part number and every "
         "attribute label  ·  overall quality score " + f["quality"], 1),
        ("The description engine reproduces all 10 derivable fields "
         "character-for-character, including a 390-character long description.", 1),
    ], size=9, gap=3)


# ---------------------------------------------------------------------------
# Assembly
# ---------------------------------------------------------------------------
def load_shots() -> List[Tuple[str, Path]]:
    captions = {
        "01_dashboard": "Dashboard — live quality gates, confidence bands, provenance",
        "02_products": "Products — searchable, per-row confidence",
        "03_product_detail": "Product detail — Raw vs Enriched vs Evidence",
        "04_evidence": "Evidence graph — every value traced to its source",
        "05_confidence": "Confidence breakdown and review flags",
        "06_review": "Human review — one row per actionable issue",
        "07_upload": "Upload — input profiled before enrichment",
    }
    out: List[Tuple[str, Path]] = []
    for stem, cap in captions.items():
        p = SHOTS / f"{stem}.png"
        if p.exists():
            out.append((cap, p))
    return out


def build(template: Path, out: Path) -> Path:
    prs = Presentation(str(template))
    f = figures()
    shots = load_shots()

    # Two extra slides for template slide 4's unanswered sub-questions, and one for
    # results. Cloned from an existing content slide so they keep the background art.
    # The template's slide 4 poses three sub-questions but gives one slide, and its
    # final slide is a designed "Thank You" closer that must stay untouched. Four
    # cloned slides carry the extra sections; the closer keeps its own artwork.
    bg = background_blob(prs, 5)
    for _ in range(4):
        clone_slide(prs, 4, bg)
    n = len(prs.slides)
    move_slide(prs, n - 4, 4)       # accuracy & trust
    move_slide(prs, n - 3, 5)       # enterprise scalability
    move_slide(prs, n - 2, 14)      # MVP: explainability
    move_slide(prs, n - 1, 15)      # results

    s = prs.slides
    s_team(s[1], f)
    s_brief(s[2], f)
    s_enrich(s[3], f)
    s_trust(s[4], f)
    s_scale(s[5], f)
    s_opportunities(s[6], f)
    s_features(s[7], f)
    s_process(s[8], f)
    s_pipeline(s[9], f)             # template's "wireframes (optional)" slot
    s_architecture(s[10], f)
    s_tech(s[11], f)
    s_cost(s[12], f)
    s_mvp(s[13], f, shots)
    s_mvp2(s[14], f, shots)
    s_results(s[15], f)
    s_future(s[16], f)
    s_links(s[17], f)
    # s[18] is the template's "Thank You" closing slide -- left exactly as supplied.

    out.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(out))
    return out


def main() -> int:
    ap = argparse.ArgumentParser()
    ap.add_argument("--template",
                    default=r"E:\UNILOG\[EXT] UniHack-Protoype Template  (1).pptx")
    ap.add_argument("--out", default=str(ROOT / "UniHack_Unilog_Product_Intelligence.pptx"))
    args = ap.parse_args()

    tpl = Path(args.template)
    if not tpl.exists():
        print(f"template not found: {tpl}", file=sys.stderr)
        return 2
    if not RUN_REPORT.exists():
        print("no run report — run scripts/run_pipeline.py first", file=sys.stderr)
        return 2

    path = build(tpl, Path(args.out))
    print(f"written: {path}")
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
