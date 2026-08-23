"""
Native PowerPoint diagrams for the UniHack deck.

Drawn with real pptx shapes and connectors rather than pasted images, so they stay
crisp at any zoom, remain editable, and inherit the template's white content area.

Three diagrams, each answering a different question a judge will ask:

    process_flow()      what happens to a row, end to end
    architecture()      how the system is layered
    ai_pipeline()       where the AI sits and what stops it inventing facts
"""
from __future__ import annotations

from typing import Dict, List, Optional, Sequence, Tuple

from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Emu, Inches, Pt

# Palette — cool technical, matching the project report.
INK = RGBColor(0x13, 0x1A, 0x21)
MUTED = RGBColor(0x5C, 0x6B, 0x75)
RULE = RGBColor(0xC5, 0xD0, 0xD6)
ACCENT = RGBColor(0x0F, 0x6E, 0x77)
ACCENT_FILL = RGBColor(0xE2, 0xF0, 0xF1)
DETERM = RGBColor(0x2F, 0x7D, 0x4F)          # deterministic stage
DETERM_FILL = RGBColor(0xE6, 0xF2, 0xEA)
AI = RGBColor(0x8A, 0x64, 0x10)              # AI-assisted stage
AI_FILL = RGBColor(0xF7, 0xEF, 0xDD)
HUMAN = RGBColor(0xA3, 0x3A, 0x32)           # human-in-the-loop
HUMAN_FILL = RGBColor(0xF8, 0xE8, 0xE6)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

FONT = "Google Sans"


# ---------------------------------------------------------------------------
# Primitives
# ---------------------------------------------------------------------------
def box(slide, left: float, top: float, width: float, height: float, text: str,
        *, fill: RGBColor = WHITE, line: RGBColor = RULE, font_color: RGBColor = INK,
        size: float = 9.0, bold: bool = False, shape=MSO_SHAPE.ROUNDED_RECTANGLE,
        line_w: float = 0.75):
    sh = slide.shapes.add_shape(shape, Inches(left), Inches(top),
                                Inches(width), Inches(height))
    sh.fill.solid()
    sh.fill.fore_color.rgb = fill
    sh.line.color.rgb = line
    sh.line.width = Pt(line_w)
    sh.shadow.inherit = False
    if hasattr(sh, "adjustments") and len(sh.adjustments):
        try:
            sh.adjustments[0] = 0.10          # subtle corner radius
        except (IndexError, ValueError):
            pass

    tf = sh.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Inches(0.04)
    tf.margin_top = tf.margin_bottom = Inches(0.02)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    lines = text.split("\n")
    for i, line_text in enumerate(lines):
        para = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        para.alignment = PP_ALIGN.CENTER
        run = para.add_run()
        run.text = line_text
        run.font.name = FONT
        run.font.size = Pt(size if i == 0 else size - 1.2)
        run.font.bold = bold and i == 0
        run.font.color.rgb = font_color if i == 0 else MUTED
    return sh


def label(slide, left: float, top: float, width: float, text: str,
          *, size: float = 8.0, color: RGBColor = MUTED, bold: bool = False,
          align=PP_ALIGN.LEFT, caps: bool = False):
    tb = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(0.24))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_top = tf.margin_bottom = 0
    para = tf.paragraphs[0]
    para.alignment = align
    run = para.add_run()
    run.text = text.upper() if caps else text
    run.font.name = FONT
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return tb


def arrow(slide, x1: float, y1: float, x2: float, y2: float,
          color: RGBColor = ACCENT, width: float = 1.25):
    c = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(x1), Inches(y1),
                                   Inches(x2), Inches(y2))
    c.line.color.rgb = color
    c.line.width = Pt(width)
    # arrowhead
    ln = c.line._get_or_add_ln()
    from pptx.oxml.ns import qn
    from lxml import etree
    tail = etree.SubElement(ln, qn("a:tailEnd"))
    tail.set("type", "triangle")
    tail.set("w", "sm")
    tail.set("len", "sm")
    return c


def legend(slide, left: float, top: float,
           items: Sequence[Tuple[str, RGBColor, RGBColor]]):
    x = left
    for text, line_c, fill_c in items:
        sw = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x),
                                    Inches(top), Inches(0.14), Inches(0.14))
        sw.fill.solid()
        sw.fill.fore_color.rgb = fill_c
        sw.line.color.rgb = line_c
        sw.line.width = Pt(0.75)
        sw.shadow.inherit = False
        label(slide, x + 0.20, top - 0.03, 1.9, text, size=7.5, color=MUTED)
        x += 0.20 + 0.12 * len(text) + 0.18


# ---------------------------------------------------------------------------
# Diagram 1 — end-to-end process flow
# ---------------------------------------------------------------------------
def process_flow(slide, top: float = 1.32) -> None:
    """One row's journey: raw feed -> 252-column record, with the two exits."""
    label(slide, 0.42, top - 0.24, 9.2,
          "one catalogue row, end to end   ·   green = deterministic   "
          "amber = AI-assisted with fallback   red = human", size=8, color=MUTED)

    y = top + 0.10
    bw, bh, gap = 1.42, 0.58, 0.13
    x0 = 0.42                       # 6*1.42 + 5*0.13 = 9.17, fits in 10in

    stages = [
        ("Profile\n& clean", DETERM, DETERM_FILL),
        ("Resolve\nentities", DETERM, DETERM_FILL),
        ("Parse\nproduct", AI, AI_FILL),
        ("Classify\ncategory", AI, AI_FILL),
        ("Retrieve\nLOV", ACCENT, ACCENT_FILL),
        ("Extract\nattributes", DETERM, DETERM_FILL),
    ]
    x = x0
    for i, (text, line_c, fill_c) in enumerate(stages):
        box(slide, x, y, bw, bh, text, fill=fill_c, line=line_c,
            font_color=line_c, size=9, bold=True)
        if i < len(stages) - 1:
            arrow(slide, x + bw, y + bh / 2, x + bw + gap, y + bh / 2)
        x += bw + gap

    # second row, right to left
    y2 = y + bh + 0.52
    stages2 = [
        ("Normalise\nfractions · UOM", DETERM, DETERM_FILL),
        ("Research\nmanufacturer", ACCENT, ACCENT_FILL),
        ("Generate\ndescriptions", DETERM, DETERM_FILL),
        ("Validate\n+ self-correct", DETERM, DETERM_FILL),
        ("Score\nconfidence", DETERM, DETERM_FILL),
        ("Route", HUMAN, HUMAN_FILL),
    ]
    x = x0
    for i, (text, line_c, fill_c) in enumerate(stages2):
        box(slide, x, y2, bw, bh, text, fill=fill_c, line=line_c,
            font_color=line_c, size=9, bold=True)
        if i < len(stages2) - 1:
            arrow(slide, x + bw, y2 + bh / 2, x + bw + gap, y2 + bh / 2)
        x += bw + gap

    # 'continues below' marker, kept inside the canvas
    last_x = x0 + 6 * (bw + gap) - gap
    label(slide, last_x - 1.75, y + bh + 0.06, 1.75, "continues below  ↴",
          size=7, color=MUTED, align=PP_ALIGN.RIGHT)

    # outcomes
    y3 = y2 + bh + 0.46
    box(slide, 3.40, y3, 2.85, 0.52,
        "PUBLISH  ·  533 rows\n252-column delivery record",
        fill=DETERM_FILL, line=DETERM, font_color=DETERM, size=9, bold=True)
    box(slide, 6.45, y3, 2.85, 0.52,
        "HUMAN REVIEW  ·  467 rows\nspecific, actionable reason each",
        fill=HUMAN_FILL, line=HUMAN, font_color=HUMAN, size=9, bold=True)
    arrow(slide, 7.60, y2 + bh, 5.5, y3, color=DETERM)
    arrow(slide, 8.30, y2 + bh, 8.30, y3, color=HUMAN)

    label(slide, 0.42, y3 + 0.10, 2.8,
          "Nothing is guessed:\nan unresolved field is left\nblank and flagged.",
          size=8, color=MUTED)


# ---------------------------------------------------------------------------
# Diagram 2 — system architecture
# ---------------------------------------------------------------------------
def architecture(slide, top: float = 1.28) -> None:
    """Layered view: interface, orchestration, agents, reference, output."""
    label(slide, 0.5, top - 0.22, 9.0,
          "every AI stage has a deterministic fallback — the system runs and reports "
          "honestly with no API key", size=8, color=MUTED)

    lanes: List[Tuple[str, float, List[Tuple[str, RGBColor, RGBColor]]]] = [
        ("INTERFACE", 0.62, [
            ("React + Tailwind dashboard\n7 pages", ACCENT, ACCENT_FILL),
            ("FastAPI\n15 endpoints", ACCENT, ACCENT_FILL),
            ("CLI\nrun · demo · benchmark", ACCENT, ACCENT_FILL),
            ("Export\nCSV · XLSX · review queue", ACCENT, ACCENT_FILL),
        ]),
        ("ORCHESTRATION", 0.56, [
            ("Two-pass corpus-aware orchestrator   ·   12 toggleable stages   ·   "
             "per-row isolation   ·   extraction cache", INK, WHITE),
        ]),
        ("AGENTS", 0.62, [
            ("Product\nParser", AI, AI_FILL),
            ("Manufacturer\nResolver", DETERM, DETERM_FILL),
            ("Brand\nResolver", DETERM, DETERM_FILL),
            ("Taxonomy\nAgent", AI, AI_FILL),
            ("Attribute\nAgent", DETERM, DETERM_FILL),
            ("Research\nAgent", ACCENT, ACCENT_FILL),
            ("Description\nAgent", DETERM, DETERM_FILL),
            ("Validation\nAgent", DETERM, DETERM_FILL),
        ]),
        ("KNOWLEDGE", 0.62, [
            ("RAG retrieval\n7 collections · 751 docs", ACCENT, ACCENT_FILL),
            ("Reference registry\nofficial → derived → seed", DETERM, DETERM_FILL),
            ("Normalisation\nfractions · UOM · casing", DETERM, DETERM_FILL),
            ("LLM client\nstrict JSON · cache · offline", AI, AI_FILL),
        ]),
    ]

    y = top + 0.06
    for lane_name, lane_h, items in lanes:
        label(slide, 0.5, y + lane_h / 2 - 0.10, 1.05, lane_name,
              size=7.5, color=MUTED, bold=True, caps=True)
        x = 1.62
        total_w = 7.95
        gap = 0.10
        bw = (total_w - gap * (len(items) - 1)) / len(items)
        for text, line_c, fill_c in items:
            box(slide, x, y, bw, lane_h, text, fill=fill_c, line=line_c,
                font_color=line_c, size=8.2, bold=True)
            x += bw + gap
        y += lane_h + 0.16

    # evidence spine
    box(slide, 1.62, y, 7.95, 0.42,
        "EVIDENCE  ·  every value carries source + transformation + confidence + "
        "validation status",
        fill=WHITE, line=INK, font_color=INK, size=8.5, bold=True)
    label(slide, 0.5, y + 0.11, 1.05, "TRACE", size=7.5, color=MUTED,
          bold=True, caps=True)


# ---------------------------------------------------------------------------
# Diagram 3 — AI / RAG enrichment pipeline
# ---------------------------------------------------------------------------
def ai_pipeline(slide, top: float = 1.26) -> None:
    """Where the model is allowed to act, and what stops it inventing facts."""
    label(slide, 0.5, top - 0.22, 9.1,
          "the model may STRUCTURE information — it may never ADD it", size=8,
          color=MUTED)

    y = top + 0.08

    # Left: input + retrieval
    box(slide, 0.5, y, 2.05, 0.54, "Raw row\n6 columns, 3 placeholders",
        fill=WHITE, line=INK, font_color=INK, size=9, bold=True)

    box(slide, 0.5, y + 0.74, 2.05, 1.50,
        "RAG COLLECTIONS\n\nmanufacturer + brand\nglobal LOV\nUOM standards\n"
        "guidelines\nexamples\nfittings · faucets",
        fill=ACCENT_FILL, line=ACCENT, font_color=ACCENT, size=8.5, bold=True)
    label(slide, 0.5, y + 2.28, 2.05, "751 documents · filter BEFORE rank",
          size=7, color=MUTED)

    # Middle: the three guarded call sites
    gx = 2.95
    gates = [
        ("1  Product parser", "fires when the row is opaque",
         "GUARD  every returned token must\nappear in the source text"),
        ("2  Taxonomy tie-break", "fires on ambiguity — 53 of 1,000 rows",
         "GUARD  must pick a supplied id or\n'none'; confidence capped at 0.72"),
        ("3  LOV rescue", "fires when a value misses the vocabulary",
         "GUARD  only an exact alias inside the\nsame attribute is applied"),
    ]
    gy = y
    for title, when, guard in gates:
        box(slide, gx, gy, 3.55, 0.42, title, fill=AI_FILL, line=AI,
            font_color=AI, size=9, bold=True)
        label(slide, gx + 0.06, gy + 0.44, 3.45, when, size=7.5, color=MUTED)
        box(slide, gx, gy + 0.64, 3.55, 0.42, guard, fill=WHITE, line=DETERM,
            font_color=DETERM, size=7.5)
        gy += 1.20

    arrow(slide, 2.55, y + 0.27, 2.95, y + 0.27)
    arrow(slide, 2.55, y + 1.4, 2.95, y + 1.4, color=ACCENT)

    # Right: outcome
    ox = 6.68
    box(slide, ox, y, 2.82, 0.80,
        "VERIFIED FACT SHEET\nresolved entities + extracted attributes only",
        fill=DETERM_FILL, line=DETERM, font_color=DETERM, size=8.5, bold=True)
    arrow(slide, 6.50, y + 0.40, ox, y + 0.40, color=DETERM)

    box(slide, ox, y + 1.00, 2.82, 0.72,
        "Description generator\nnever sees raw text",
        fill=WHITE, line=INK, font_color=INK, size=8.5, bold=True)
    arrow(slide, ox + 1.41, y + 0.80, ox + 1.41, y + 1.00, color=DETERM)

    box(slide, ox, y + 1.92, 2.82, 0.72,
        "Validate → self-correct\nconfidence → review",
        fill=HUMAN_FILL, line=HUMAN, font_color=HUMAN, size=8.5, bold=True)
    arrow(slide, ox + 1.41, y + 1.72, ox + 1.41, y + 1.92, color=HUMAN)

    label(slide, ox, y + 2.70, 2.82,
          "Result: the generator is structurally\nincapable of inventing a fact.",
          size=7.5, color=MUTED)
